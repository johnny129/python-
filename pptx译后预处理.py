import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Pt

def adjust_text_format(presentation, font_scale, line_spacing, apply_spacing):
    """
    调整 PowerPoint 演示文稿的文本格式，包括字体缩放和行距。

    参数:
        presentation (Presentation): PowerPoint 演示文稿对象。
        font_scale (float): 字体大小的缩放比例。
        line_spacing (float): 应用的行距。
        apply_spacing (bool): 是否应用行距调整。

    返回:
        Presentation: 调整后的 PowerPoint 演示文稿对象。
    """
    def adjust_shape_text(shape):
        """调整演示文稿中形状内的文本，包括文本框和表格单元格。"""
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # 为空行添加空格以避免空段落问题
                if not paragraph.text.strip():
                    paragraph.text = " "
                    # 设置新添加空格的字体大小
                    for run in paragraph.runs:
                        if run.font.size is None:
                            run.font.size = Pt(10 * font_scale)  # 假设默认字体大小为10
                        else:
                            run.font.size = Pt(run.font.size.pt * font_scales)
                
                for run in paragraph.runs:
                    # 根据提供的缩放比例调整字体大小
                    if run.font.size is not None:
                        run.font.size = Pt(run.font.size.pt * font_scale)
                # 如果设置了行距标志，则设置段落行距
                if apply_spacing:
                    paragraph.space_after = 0
                    paragraph.space_before = 0
                    paragraph.line_spacing = line_spacing

        # 如果形状包含表格，则调整每个单元格中的文本
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:  # 确保单元格具有 text_frame
                        for paragraph in cell.text_frame.paragraphs:
                            # 为空行添加空格以避免空段落问题
                            if not paragraph.text.strip():
                                paragraph.text = " "
                                # 设置新添加空格的字体大小
                                for run in paragraph.runs:
                                    if run.font.size is None:
                                        run.font.size = Pt(18 * font_scale)  # 假设默认字体大小为18
                                    else:
                                        run.font.size = Pt(run.font.size.pt * font_scale)
                            
                            for run in paragraph.runs:
                                # 根据提供的缩放比例调整字体大小
                                if run.font.size is not None:
                                    run.font.size = Pt(run.font.size.pt * font_scale)
                            # 如果设置了行距标志，则设置段落行距
                            if apply_spacing:
                                paragraph.space_after = 0
                                paragraph.space_before = 0
                                paragraph.line_spacing = line_spacing

        # 如果形状是组合形状，则递归调整子形状中的文本
        if shape.shape_type == 6:  # 6 表示组合形状
            for sub_shape in shape.shapes:
                adjust_shape_text(sub_shape)

    # 遍历演示文稿中的所有幻灯片并调整每个形状中的文本
    for slide in presentation.slides:
        for shape in slide.shapes:
            adjust_shape_text(shape)

    # 遍历所有幻灯片母版并调整每个形状中的文本
    for slide_master in presentation.slide_masters:
        for layout in slide_master.slide_layouts:
            for shape in layout.shapes:
                adjust_shape_text(shape)

    return presentation

def process_ppt(input_path, output_folder, result_text, font_scale, line_spacing, apply_spacing):
    """
    通过调整文本格式处理单个 PPT 文件。

    参数:
        input_path (str): 输入 PPT 文件的文件路径。
        output_folder (str): 保存调整后的 PPT 文件的文件夹。
        result_text (Text): 显示结果的 Tkinter Text 小部件。
        font_scale (float): 字体大小的缩放比例。
        line_spacing (float): 应用的行距。
        apply_spacing (bool): 是否应用行距调整。
    """
    try:
        # 从输入文件路径加载演示文稿
        presentation = Presentation(input_path)
        # 调整演示文稿的文本格式
        adjusted_presentation = adjust_text_format(presentation, font_scale, line_spacing, apply_spacing)

        # 如果输出文件夹不存在，则创建它
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # 将调整后的演示文稿保存到输出文件夹
        output_ppt = os.path.join(output_folder, os.path.basename(input_path))
        adjusted_presentation.save(output_ppt)

        # 在结果 Text 小部件中显示成功消息
        result_text.config(state=tk.NORMAL)
        result_text.insert(tk.END, f"处理成功：{output_ppt}\n")
        result_text.config(state=tk.DISABLED)
    except Exception as e:
        # 在结果 Text 小部件中显示错误消息
        result_text.config(state=tk.NORMAL)
        result_text.insert(tk.END, f"处理失败：{os.path.basename(input_path)}，错误: {str(e)}\n")
        result_text.config(state=tk.DISABLED)

def browse_folder(entry):
    """提示用户选择文件夹并在提供的 Entry 小部件中显示路径。"""
    selected_path = filedialog.askdirectory()
    entry.delete(0, tk.END)
    entry.insert(0, selected_path)

def browse_file(entry):
    """提示用户选择单个文件并在提供的 Entry 小部件中显示路径。"""
    selected_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    entry.delete(0, tk.END)
    entry.insert(0, selected_path)

def process():
    """处理按钮的回调函数，用于处理 PPT 文件。"""
    folder_path = entry_folder.get()
    file_path = entry_file.get()
    result_text.config(state=tk.NORMAL)
    result_text.delete(1.0, tk.END)

    # 获取用户输入的字体缩放比例和行距
    font_scale = float(entry_font_scale.get())
    line_spacing = float(entry_line_spacing.get())
    apply_spacing = apply_spacing_var.get()

    # 处理选定文件夹中的所有 PPT 文件或单个选定文件
    if folder_path:
        output_folder = os.path.join(folder_path, "output")
        for filename in os.listdir(folder_path):
            if filename.endswith(".pptx"):
                input_path = os.path.join(folder_path, filename)
                process_ppt(input_path, output_folder, result_text, font_scale, line_spacing, apply_spacing)
    elif file_path:
        output_folder = os.path.join(os.path.dirname(file_path), "output")
        process_ppt(file_path, output_folder, result_text, font_scale, line_spacing, apply_spacing)

    result_text.config(state=tk.DISABLED)

# 创建主 Tkinter 窗口
root = tk.Tk()
root.title("PPT 译后预处理工具")

# 创建 UI 组件
label_folder = tk.Label(root, text="选择文件夹:")
entry_folder = tk.Entry(root, width=50)
btn_browse_folder = tk.Button(root, text="选择文件夹", bg="#2cc2d9", fg="#FFFFFF", command=lambda: browse_folder(entry_folder))

label_file = tk.Label(root, text="选择单文件:")
entry_file = tk.Entry(root, width=50)
btn_browse_file = tk.Button(root, text="选择单文件", bg="#00da6a", fg="#FFFFFF", command=lambda: browse_file(entry_file))

label_font_scale = tk.Label(root, text="字体缩放比例:")
entry_font_scale = tk.Entry(root, width=10)
entry_font_scale.insert(0, "0.6")

label_line_spacing = tk.Label(root, text="行距设置:")
entry_line_spacing = tk.Entry(root, width=10)
entry_line_spacing.insert(0, "1.0")

apply_spacing_var = tk.BooleanVar()
chk_apply_spacing = tk.Checkbutton(root, text="修改行距", fg="#f01363", variable=apply_spacing_var)

btn_process = tk.Button(root, text="处理", bg="#b80001", fg="#FFFFFF", command=process)

result_text = tk.Text(root, wrap=tk.WORD, height=20, width=50, state=tk.DISABLED, bg="white")

# 使用网格布局排列 UI 组件
label_folder.grid(row=0, column=0, pady=5, sticky="w")
entry_folder.grid(row=0, column=1, pady=5, sticky="ew")
btn_browse_folder.grid(row=0, column=2, padx=5, pady=5)

label_file.grid(row=1, column=0, pady=5, sticky="w")
entry_file.grid(row=1, column=1, pady=5, sticky="ew")
btn_browse_file.grid(row=1, column=2, padx=5, pady=5)

label_font_scale.grid(row=2, column=0, pady=5, sticky="w")
entry_font_scale.grid(row=2, column=1, pady=5, sticky="w")

label_line_spacing.grid(row=3, column=0, pady=5, sticky="w")
entry_line_spacing.grid(row=3, column=1, pady=5, sticky="w")

chk_apply_spacing.grid(row=5, column=0, pady=5, sticky="w")

btn_process.grid(row=6, column=0, columnspan=3, pady=10)

result_text.grid(row=7, column=0, columnspan=3, padx=10, pady=10, sticky="nsew")

# 调整窗口比例
root.columnconfigure(1, weight=1)
root.rowconfigure(7, weight=1)

# 启动 Tkinter 主循环
root.mainloop()

