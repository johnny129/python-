import sys
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images_from_ppt(ppt_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    presentation = Presentation(ppt_path)
    image_count = 0

    for slide_index, slide in enumerate(presentation.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext

                image_count += 1
                image_filename = f"slide_{slide_index+1}_image_{image_count}.{image_format}"
                image_path = os.path.join(output_dir, image_filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

    print(f"图片提取完成！共提取 {image_count} 张图片。")

if __name__ == "__main__":
    ppt_file = sys.argv[1]
    output_dir = os.path.join(os.path.dirname(ppt_file), os.path.splitext(os.path.basename(ppt_file))[0])
    extract_images_from_ppt(ppt_file, output_dir)

