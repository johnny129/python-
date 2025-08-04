import os
import json
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinterdnd2 import TkinterDnD, DND_FILES
import winreg as reg
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile

# 配置文件路径
CONFIG_FILE = 'config.json'

def load_config():
    if os.path.exists(CONFIG_FILE):
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)
            return config
    return {}

def save_config(config):
    with open(CONFIG_FILE, 'w') as f:
        json.dump(config, f)

# 检查注册表中是否存在右键菜单项
def check_context_menu_exists():
    try:
        with reg.OpenKey(reg.HKEY_CLASSES_ROOT, r'SystemFileAssociations\.pptx\shell\ExtractImages', 0, reg.KEY_READ) as key:
            return True
    except FileNotFoundError:
        return False

# 生成注册表文件
def generate_reg_file(action):
    exe_path = os.path.abspath("extract_images_cmd.exe").replace("\\", "\\\\")
    reg_content = ""

    if action == 'add':
        reg_content = f"""Windows Registry Editor Version 5.00

[HKEY_CLASSES_ROOT\\SystemFileAssociations\\.pptx\\shell\\ExtractImages]
@="提取PPT图片"
"Icon"="C:\\\Windows\\\System32\\\shell32.dll,46"

[HKEY_CLASSES_ROOT\\SystemFileAssociations\\.pptx\\shell\\ExtractImages\\command]
@="\\\"{exe_path}\\\" \\"%1\\\""
"""
    elif action == 'delete':
        reg_content = """Windows Registry Editor Version 5.00

[-HKEY_CLASSES_ROOT\\SystemFileAssociations\\.pptx\\shell\\ExtractImages]
"""

    # 将 .reg 文件保存在 %temp% 目录
    reg_file_path = os.path.join(tempfile.gettempdir(), 'context_menu.reg')
    with open(reg_file_path, 'w') as f:
        f.write(reg_content)
    
    return reg_file_path

# 添加右键菜单
def add_to_context_menu():
    reg_file_path = generate_reg_file('add')
    os.startfile(reg_file_path)  # 打开注册表文件以供用户双击导入

# 删除右键菜单
def remove_from_context_menu():
    reg_file_path = generate_reg_file('delete')
    os.startfile(reg_file_path)  # 打开注册表文件以供用户双击导入

# 切换右键菜单状态
def toggle_context_menu():
    if context_menu_var.get():
        add_to_context_menu()
    else:
        remove_from_context_menu()

# 选择PPT文件
def select_ppt_file():
    ppt_file_path = filedialog.askopenfilename(filetypes=[("PPT files", "*.pptx")])
    if ppt_file_path:
        ppt_file_path = os.path.normpath(ppt_file_path)  # 标准化路径
        ppt_entry.delete(0, tk.END)
        ppt_entry.insert(0, ppt_file_path)
        default_output_dir.set(os.path.join(os.path.dirname(ppt_file_path), os.path.splitext(os.path.basename(ppt_file_path))[0]))

# 选择保存路径
def select_output_dir():
    ppt_file_path = ppt_entry.get()
    output_dir = filedialog.askdirectory()

    if output_dir:
        output_dir = os.path.normpath(output_dir)  # 标准化路径
        output_dir_entry.delete(0, tk.END)
        ppt_name = os.path.splitext(os.path.basename(ppt_file_path))[0]
        output_dir_entry.insert(0, os.path.join(output_dir, ppt_name))

# 处理拖放
def drop(event):
    ppt_file_path = event.data.strip('{}')  # 去掉花括号
    ppt_file_path = os.path.normpath(ppt_file_path)  # 标准化路径
    ppt_entry.delete(0, tk.END)
    ppt_entry.insert(0, ppt_file_path)
    default_output_dir.set(os.path.join(os.path.dirname(ppt_file_path), os.path.splitext(os.path.basename(ppt_file_path))[0]))

# 提取图片
def extract_images():
    ppt_file_path = os.path.normpath(ppt_entry.get())  # 标准化路径
    output_dir = os.path.normpath(output_dir_entry.get())  # 标准化路径

    if not ppt_file_path or not output_dir:
        messagebox.showerror("错误", "请指定PPT文件和保存路径！")
        return

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    presentation = Presentation(ppt_file_path)
    image_count = 0

    for slide_index, slide in enumerate(presentation.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext

                image_count += 1
                image_filename = f"slide_{slide_index+1}_image_{image_count}.{image_format}"
                image_path = os.path.normpath(os.path.join(output_dir, image_filename))  # 标准化路径

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

    # 提示是否打开保存目录
    if messagebox.askyesno("提取完成", f"共提取 {image_count} 张图片！是否打开保存目录？"):
        os.startfile(output_dir)
        
# 退出时删除临时的 .reg 文件
def on_exit():
    reg_file_path = os.path.join(tempfile.gettempdir(), 'context_menu.reg')
    if os.path.exists(reg_file_path):
        os.remove(reg_file_path)
    root.destroy()

# 创建主窗口
root = TkinterDnD.Tk()
root.title("PPT图片提取工具")
root.protocol("WM_DELETE_WINDOW", on_exit)

# 整个窗口响应拖放文件
root.drop_target_register(DND_FILES)
root.dnd_bind('<<Drop>>', drop)

# PPT文件路径输入框和选择按钮
ppt_label = tk.Label(root, text="选择PPT文件:")
ppt_label.grid(row=0, column=0, padx=10, pady=5, sticky="e")
ppt_entry = tk.Entry(root, width=50)
ppt_entry.grid(row=0, column=1, padx=10, pady=5)
ppt_button = tk.Button(root, text="浏览", command=select_ppt_file)
ppt_button.grid(row=0, column=2, padx=10, pady=5)

# 保存路径输入框和选择按钮
output_dir_label = tk.Label(root, text="选择保存路径:")
output_dir_label.grid(row=1, column=0, padx=10, pady=5, sticky="e")
default_output_dir = tk.StringVar()
output_dir_entry = tk.Entry(root, width=50, textvariable=default_output_dir)
output_dir_entry.grid(row=1, column=1, padx=10, pady=5)
output_dir_button = tk.Button(root, text="浏览", command=select_output_dir)
output_dir_button.grid(row=1, column=2, padx=10, pady=5)

# 提取图片按钮
extract_button = tk.Button(root, text="提取图片", command=extract_images)
extract_button.grid(row=2, column=1, padx=10, pady=20)

# 右键菜单选项
context_menu_var = tk.IntVar()
context_menu_checkbox = tk.Checkbutton(root, text="添加到右键菜单", variable=context_menu_var, command=toggle_context_menu)
context_menu_checkbox.grid(row=3, column=1, padx=10, pady=5)

# 启动时检查并设置右键菜单状态
if check_context_menu_exists():
    context_menu_var.set(1)

# 运行主循环
root.mainloop()
