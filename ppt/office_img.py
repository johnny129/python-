import os
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinterdnd2 import TkinterDnD, DND_FILES
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import zipfile
from PyPDF2 import PdfReader
import pythoncom
from win32com.client import Dispatch

def select_file():
    file_path = filedialog.askopenfilename(filetypes=[("Office and PDF files", "*.pptx *.docx *.xlsx *.pdf *.doc *.xls *.ppt")])
    if file_path:
        file_folder_entry.delete(0, tk.END)
        file_folder_entry.insert(0, file_path)

def select_folder():
    folder_path = filedialog.askdirectory()
    if folder_path:
        file_folder_entry.delete(0, tk.END)
        file_folder_entry.insert(0, folder_path)

def select_output_dir():
    output_dir = filedialog.askdirectory()
    if output_dir:
        output_dir_entry.delete(0, tk.END)
        output_dir_entry.insert(0, output_dir)

def get_supported_files(folder, selected_types):
    files = []
    for root, _, fs in os.walk(folder):
        for f in fs:
            file_ext = os.path.splitext(f)[1].lower()
            if file_ext in selected_types:
                files.append(os.path.join(root, f))
    return files

def convert_to_office_new_format(file_path):
    pythoncom.CoInitialize()
    ext = os.path.splitext(file_path)[1].lower()
    new_path = file_path + 'x'
    try:
        if ext == '.doc':
            word = Dispatch('Word.Application')
            word.Visible = False
            doc = word.Documents.Open(file_path)
            doc.SaveAs(new_path, FileFormat=12)
            doc.Close()
            word.Quit()
        elif ext == '.xls':
            excel = Dispatch('Excel.Application')
            excel.Visible = False
            wb = excel.Workbooks.Open(file_path)
            wb.SaveAs(new_path, FileFormat=51)
            wb.Close()
            excel.Quit()
        elif ext == '.ppt':
            powerpoint = Dispatch('PowerPoint.Application')
            powerpoint.Visible = False
            ppt = powerpoint.Presentations.Open(file_path, WithWindow=False)
            ppt.SaveAs(new_path, FileFormat=24)
            ppt.Close()
            powerpoint.Quit()
        else:
            return None
        return new_path
    except Exception:
        return None

def extract_images_from_ppt(ppt_path, output_dir):
    presentation = Presentation(ppt_path)
    image_count = 0
    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext
                image_count += 1
                image_filename = f"slide_{slide_index+1}_image_{image_count}.{image_format}"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
    return image_count

def extract_images_from_word(word_path, output_dir):
    count = 0
    with zipfile.ZipFile(word_path) as docx_zip:
        for name in docx_zip.namelist():
            if name.startswith("word/media/"):
                img_data = docx_zip.read(name)
                img_filename = os.path.join(output_dir, os.path.basename(name))
                with open(img_filename, "wb") as f:
                    f.write(img_data)
                count += 1
    return count

def extract_images_from_excel(excel_path, output_dir):
    count = 0
    with zipfile.ZipFile(excel_path) as xlsx_zip:
        for name in xlsx_zip.namelist():
            if name.startswith("xl/media/"):
                img_data = xlsx_zip.read(name)
                img_filename = os.path.join(output_dir, os.path.basename(name))
                with open(img_filename, "wb") as f:
                    f.write(img_data)
                count += 1
    return count

def extract_images_from_pdf(pdf_path, output_dir):
    reader = PdfReader(pdf_path)
    image_count = 0
    for page_index, page in enumerate(reader.pages):
        if '/XObject' in page['/Resources']:
            xObject = page['/Resources']['/XObject'].get_object()
            for obj in xObject:
                if xObject[obj]['/Subtype'] == '/Image':
                    image_data = xObject[obj].get_data()
                    if '/Filter' in xObject[obj]:
                        filter_name = xObject[obj]['/Filter']
                        if filter_name == '/DCTDecode':
                            image_format = 'jpg'
                        elif filter_name == '/JPXDecode':
                            image_format = 'jp2'
                        elif filter_name == '/FlateDecode':
                            image_format = 'png'
                        else:
                            image_format = 'jpg'
                    else:
                        image_format = 'bin'
                    image_count += 1
                    image_filename = f"page_{page_index+1}_image_{image_count}.{image_format}"
                    image_path = os.path.join(output_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_data)
    return image_count

def process_files(file_paths, output_root):
    results = []
    for file_path in file_paths:
        ext = os.path.splitext(file_path)[1].lower()
        file_name = os.path.splitext(os.path.basename(file_path))[0]
        output_dir = os.path.join(output_root, f"图片-{file_name}")
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        if ext in ['.doc', '.xls', '.ppt']:
            new_file_path = convert_to_office_new_format(file_path)
            if new_file_path and os.path.exists(new_file_path):
                file_path = new_file_path
                ext = os.path.splitext(file_path)[1].lower()
            else:
                results.append((file_path, 0))
                continue
        if ext == '.pptx':
            count = extract_images_from_ppt(file_path, output_dir)
            results.append((file_path, count))
        elif ext == '.docx':
            count = extract_images_from_word(file_path, output_dir)
            results.append((file_path, count))
        elif ext == '.xlsx':
            count = extract_images_from_excel(file_path, output_dir)
            results.append((file_path, count))
        elif ext == '.pdf':
            count = extract_images_from_pdf(file_path, output_dir)
            results.append((file_path, count))
    return results

def extract_images():
    output_root = output_dir_entry.get().strip() or os.getcwd()
    file_or_folder = file_folder_entry.get().strip()
    selected_types = []
    if word_var.get():
        selected_types.extend(['.docx', '.doc'])
    if excel_var.get():
        selected_types.extend(['.xlsx', '.xls'])
    if ppt_var.get():
        selected_types.extend(['.pptx', '.ppt'])
    if pdf_var.get():
        selected_types.append('.pdf')
    file_paths = []
    if os.path.isfile(file_or_folder):
        if os.path.splitext(file_or_folder)[1].lower() in selected_types:
            file_paths = [os.path.normpath(file_or_folder)]
    elif os.path.isdir(file_or_folder):
        file_paths = get_supported_files(os.path.normpath(file_or_folder), selected_types)
        if not file_paths:
            messagebox.showerror("错误", "文件夹中没有可处理的文件！")
            return
    else:
        messagebox.showerror("错误", "请选择有效的文件或文件夹！")
        return
    results = process_files(file_paths, output_root)
    total = sum([c for _, c in results])
    msg = "\n".join([f"{os.path.basename(fp)}: {c} 张图片" for fp, c in results])
    if messagebox.askyesno("提取完成", f"共处理 {len(results)} 个文件，提取 {total} 张图片。\n\n详细:\n{msg}\n\n是否打开保存目录？"):
        os.startfile(output_root)

root = TkinterDnD.Tk()
root.title("Office和PDF图片批量提取工具")
root.drop_target_register(DND_FILES)
try:
    root.iconbitmap('icon.ico')
except Exception:
    pass

file_folder_label = tk.Label(root, text="选择文件或文件夹:")
file_folder_label.grid(row=0, column=0, padx=10, pady=5, sticky="e")
file_folder_entry = tk.Entry(root, width=50)
file_folder_entry.grid(row=0, column=1, padx=10, pady=5)
file_button = tk.Button(root, text="选择文件", command=select_file)
file_button.grid(row=0, column=2, padx=5, pady=5)
folder_button = tk.Button(root, text="选择文件夹", command=select_folder)
folder_button.grid(row=0, column=3, padx=5, pady=5)

file_type_label = tk.Label(root, text="筛选文件类型:")
file_type_label.grid(row=1, column=0, padx=10, pady=5, sticky="ne")
word_var = tk.BooleanVar(value=True)
excel_var = tk.BooleanVar(value=True)
ppt_var = tk.BooleanVar(value=True)
pdf_var = tk.BooleanVar(value=True)
word_checkbox = tk.Checkbutton(root, text="Word", variable=word_var)
word_checkbox.grid(row=1, column=1, padx=5, pady=2, sticky="w")
excel_checkbox = tk.Checkbutton(root, text="Excel", variable=excel_var)
excel_checkbox.grid(row=2, column=1, padx=5, pady=2, sticky="w")
ppt_checkbox = tk.Checkbutton(root, text="PPT", variable=ppt_var)
ppt_checkbox.grid(row=3, column=1, padx=5, pady=2, sticky="w")
pdf_checkbox = tk.Checkbutton(root, text="PDF", variable=pdf_var)
pdf_checkbox.grid(row=4, column=1, padx=5, pady=2, sticky="w")

output_dir_label = tk.Label(root, text="导出路径:")
output_dir_label.grid(row=5, column=0, padx=10, pady=5, sticky="e")
output_dir_entry = tk.Entry(root, width=50)
output_dir_entry.insert(0, os.getcwd())
output_dir_entry.grid(row=5, column=1, padx=10, pady=5)
output_dir_button = tk.Button(root, text="浏览", command=select_output_dir)
output_dir_button.grid(row=5, column=2, padx=10, pady=5)

extract_button = tk.Button(root, text="批量提取图片", command=extract_images)
extract_button.grid(row=6, column=1, padx=10, pady=20)

root.mainloop()
